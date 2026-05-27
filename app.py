"""
EcoMind AI - Web Version
Backend Flask
"""

from flask import Flask, render_template, request, jsonify, send_file
from flask_cors import CORS
from groq import Groq
import os
from fpdf import FPDF
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from PIL import Image, ImageDraw
import requests
import io
from datetime import datetime

app = Flask(__name__)
CORS(app)

GROQ_API_KEY = os.environ.get("GROQ_API_KEY", "")
if not GROQ_API_KEY:
    raise RuntimeError("GROQ_API_KEY no está configurada como variable de entorno")
client = Groq(api_key=GROQ_API_KEY)

# Crear directorio para archivos generados
UPLOAD_DIR = os.path.join(os.path.dirname(__file__), 'static', 'uploads')
os.makedirs(UPLOAD_DIR, exist_ok=True)

# Variables globales para sesiones
sesiones = {}

def generar_imagen_ecologica():
    """Genera una imagen ecológica con PIL"""
    path = os.path.join(UPLOAD_DIR, "eco_logo.png")
    if os.path.exists(path):
        return path
    
    img = Image.new("RGB", (400, 200), "#1B5E20")
    draw = ImageDraw.Draw(img)
    draw.ellipse([150, 20, 250, 120], fill="#4CAF50", outline="#A5D6A7", width=3)
    draw.polygon([(200, 30), (170, 100), (230, 100)], fill="#66BB6A")
    draw.rectangle([100, 120, 300, 140], fill="#8D6E63")
    draw.rectangle([180, 80, 220, 140], fill="#5D4037")
    draw.text((130, 155), "EcoMind AI", fill="#A5D6A7")
    img.save(path)
    return path

def preguntar_ia(pregunta):
    """Consulta a la IA de Groq"""
    try:
        response = client.chat.completions.create(
            model="llama-3.3-70b-versatile",
            messages=[
                {"role": "system", "content": "Eres un experto en ecología y medio ambiente."},
                {"role": "user", "content": pregunta}
            ]
        )
        return response.choices[0].message.content
    except Exception as e:
        return f"Error: {str(e)}"

@app.route('/')
def index():
    return render_template('index.html')

@app.route('/api/pregunta', methods=['POST'])
def api_pregunta():
    """API para enviar preguntas a la IA"""
    data = request.json
    pregunta = data.get('pregunta', '')
    
    if not pregunta:
        return jsonify({'error': 'Pregunta vacía'}), 400
    
    respuesta = preguntar_ia(pregunta)
    return jsonify({'respuesta': respuesta})

@app.route('/api/generar-pdf', methods=['POST'])
def api_generar_pdf():
    """API para generar PDF"""
    data = request.json
    texto = data.get('texto', '')
    
    if not texto:
        return jsonify({'error': 'Texto vacío'}), 400
    
    try:
        pdf = FPDF()
        pdf.add_page()
        pdf.set_fill_color(27, 94, 32)
        pdf.rect(0, 0, 210, 30, "F")
        pdf.set_text_color(255, 255, 255)
        pdf.set_font("Arial", "B", 18)
        pdf.set_xy(10, 8)
        pdf.cell(0, 12, "EcoMind AI - Reporte Ecologico", ln=True, align="C")
        
        pdf.set_draw_color(76, 175, 80)
        pdf.set_line_width(0.8)
        pdf.line(10, 60, 200, 60)
        
        pdf.set_text_color(30, 30, 30)
        pdf.set_font("Arial", "", 11)
        pdf.set_xy(10, 65)
        
        for linea in texto.split("\n"):
            if linea.strip():
                limpio = linea.encode("latin-1", "replace").decode("latin-1")
                pdf.multi_cell(0, 6, limpio)
        
        filename = f"reporte_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pdf"
        filepath = os.path.join(UPLOAD_DIR, filename)
        pdf.output(filepath)
        
        return jsonify({'filename': filename, 'url': f'/static/uploads/{filename}'})
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/api/generar-ppt', methods=['POST'])
def api_generar_ppt():
    """API para generar PowerPoint"""
    data = request.json
    texto = data.get('texto', '')
    
    if not texto:
        return jsonify({'error': 'Texto vacío'}), 400
    
    try:
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        
        verde_oscuro = RGBColor(27, 94, 32)
        verde_claro = RGBColor(76, 175, 80)
        blanco = RGBColor(255, 255, 255)
        
        # Slide 1
        s1 = prs.slides.add_slide(prs.slide_layouts[6])
        bg = s1.background
        bg.fill.solid()
        bg.fill.fore_color.rgb = verde_oscuro
        
        txBox = s1.shapes.add_textbox(Inches(0.5), Inches(2), Inches(12), Inches(2))
        t = txBox.text_frame.paragraphs[0]
        t.text = "EcoMind AI"
        t.font.size = Pt(54)
        t.font.color.rgb = blanco
        t.font.bold = True
        
        p = txBox.text_frame.add_paragraph()
        p.text = "Reporte de Inteligencia Artificial Ecológica"
        p.font.size = Pt(24)
        p.font.color.rgb = verde_claro
        
        # Slide 2
        s2 = prs.slides.add_slide(prs.slide_layouts[6])
        bg2 = s2.background
        bg2.fill.solid()
        bg2.fill.fore_color.rgb = RGBColor(245, 245, 245)
        
        txBox2 = s2.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(1))
        t2 = txBox2.text_frame.paragraphs[0]
        t2.text = "Contenido"
        t2.font.size = Pt(36)
        t2.font.color.rgb = verde_oscuro
        t2.font.bold = True
        
        tb = s2.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(5))
        tbf = tb.text_frame
        tbf.word_wrap = True
        
        for i, linea in enumerate(texto.split("\n")[:15]):
            if linea.strip():
                pp = tbf.paragraphs[0] if i == 0 else tbf.add_paragraph()
                pp.text = linea[:120]
                pp.font.size = Pt(16)
                pp.font.color.rgb = RGBColor(80, 80, 80)
        
        filename = f"reporte_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
        filepath = os.path.join(UPLOAD_DIR, filename)
        prs.save(filepath)
        
        return jsonify({'filename': filename, 'url': f'/static/uploads/{filename}'})
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/api/generar-txt', methods=['POST'])
def api_generar_txt():
    """API para generar TXT"""
    data = request.json
    texto = data.get('texto', '')
    
    if not texto:
        return jsonify({'error': 'Texto vacío'}), 400
    
    try:
        contenido = "="*60 + "\n"
        contenido += "  EcoMind AI - Reporte Ecologico\n"
        contenido += "="*60 + "\n\n"
        contenido += texto
        contenido += "\n\n" + "-"*60 + "\nGenerado por EcoMind AI\n"
        
        filename = f"reporte_{datetime.now().strftime('%Y%m%d_%H%M%S')}.txt"
        filepath = os.path.join(UPLOAD_DIR, filename)
        
        with open(filepath, 'w', encoding='utf-8') as f:
            f.write(contenido)
        
        return jsonify({'filename': filename, 'url': f'/static/uploads/{filename}'})
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/api/generar-imagen-ia', methods=['POST'])
def api_generar_imagen():
    """API para generar imagen con Pollinations AI"""
    data = request.json
    consulta = data.get('consulta', '')
    
    if not consulta:
        return jsonify({'error': 'Consulta vacía'}), 400
    
    try:
        prompt = f"{consulta} ecologia naturaleza sostenible"
        url = f"https://image.pollinations.ai/prompt/{prompt.replace(' ', '%20')}"
        
        response = requests.get(url, timeout=60)
        
        filename = f"imagen_{datetime.now().strftime('%Y%m%d_%H%M%S')}.png"
        filepath = os.path.join(UPLOAD_DIR, filename)
        
        with open(filepath, 'wb') as f:
            f.write(response.content)
        
        return jsonify({'filename': filename, 'url': f'/static/uploads/{filename}'})
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/static/uploads/<filename>')
def download_file(filename):
    """Descargar archivo generado"""
    filepath = os.path.join(UPLOAD_DIR, filename)
    if os.path.exists(filepath):
        return send_file(filepath, as_attachment=True)
    return jsonify({'error': 'Archivo no encontrado'}), 404

@app.errorhandler(404)
def not_found(error):
    return jsonify({'error': 'No encontrado'}), 404

@app.errorhandler(500)
def server_error(error):
    return jsonify({'error': 'Error del servidor'}), 500

if __name__ == '__main__':
    app.run(debug=True, port=5000)
